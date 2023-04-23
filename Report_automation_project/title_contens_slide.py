import collections
import collections.abc
from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표 등을 그리기 위해

prs = Presentation() # 파워풔인트 객체 선언

bullet_slide_layout = prs.slide_layouts[1] # 1: 제목 및 내용 슬라이드
slide = prs.slides.add_slide(bullet_slide_layout) # 기존에 있던 슬라이드에 추가

# 제목
title_shape = slide.placeholders[0]
title_shape.text = 'Adding a Bullet Slide'

# 내용
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

# 단락 추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2 # 2: 들여쓰기 레벨

# 저장
prs.save('test.pptx')