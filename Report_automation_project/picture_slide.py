import collections
import collections.abc
from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표 등을 그리기 위해

prs = Presentation()

img_path = 'res/powerpoint_handling/slide_test.jpg'

blank_slide_layout = prs.slide_layouts[6] # 6: 제목/내용이 없는 '빈' 슬라이드
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
width = height = Inches(1)
# width, height가 없을 경우 원본 사이즈로
pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)

left = Inches(3)
width = Inches(5.5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)

prs.save('test.pptx')